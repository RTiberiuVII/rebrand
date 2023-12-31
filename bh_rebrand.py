"""
Author: Tiberiu Rociu, Valentin Kiss, Ole Gerlof
Date: 21/06/2023
Version: 1.3
Purpose:
Rebrand MS Office documents (Word, Excel PowerPoint) from
    former â€œBHGEâ€ to current Baker Hughes Company style and re-create PDF files
Applies to company logo, company name, document font, brand colors

"""
import multiprocessing
import sys
import os
import xml
from time import time, strftime, gmtime, sleep
from datetime import datetime
from re import search, findall
from zipfile import ZipFile, ZIP_DEFLATED

import lxml.etree
import pptx
from docx import Document
import comtypes.client
import xml.etree.ElementTree as ET
import win32com.client as win32
import shutil
import pathlib
from PIL import Image
import numpy as np
from pptx import Presentation
from pdf2docx import Converter as ConverterPdf2Docx
from docx2pdf import convert as ConvertDocx2Pdf

FILE_FORMAT_PDF_WORD = 17
FILE_FORMAT_PDF_EXCEL = 0
FILE_FORMAT_PDF_PPT = 2
FILE_FORMAT_XLSX = 51
FILE_FORMAT_DOCX = 16
different_logos_found = 800
image_comparisons = 0
header_images = {}
file_run_times = {
    'excel': 0.0,
    'word': 0.0,
    'powerpoint': 0.0,
    'pdf': 0.0
}

log_information = {
    'Inputfile': '',
    'Logo': 'No Logo Found',
    'NumLogos': 0,
    'Notes': '',
    'LegacyText': '',
    'Time': '',
    'Warning': ''
}

FOLDERS = ['OutputFolder', 'LogFolder', 'BetweenFolder', 'HeaderImageReplacedFoler', 'FoundLogosFolder', 'ImagesFolder']

# Hold PDF Document ID numbers to keep track of .docx files that need to be converted back to PDF
pdfs = []

"""def count_docx(file_name):
    document = Document(file_name)

    newparatextlist = []
    for paratext in document.paragraphs:
        newparatextlist.append(paratext.text)

    return len(re.findall(r'\w+', '\n'.join(newparatextlist)))


def shuttle_text(shuttle):
    t = ''
    for i in shuttle:
        t += i.text
    return t"""


def replace_text(runs, target, replace):
    """
    Replaces the target string with the replace string in the runs.

    Parameters:
    - runs (List[Run]): List of runs to search and replace within.
    - target (str): The string to be replaced.
    - replace (str): The replacement string.
    """
    begin = 0
    full_text = ""
    for end, run in enumerate(runs):
        if run.text is not None:
            full_text += run.text

            if target in full_text:

                # Find the beginning index
                index = full_text.index(target)
                while index >= len(runs[begin].text):
                    index -= len(runs[begin].text)
                    begin += 1

                # Find the corresponding runs
                shuttle = runs[begin:end + 1]

                # Perform the replace operation
                if target in shuttle[0].text:
                    shuttle[0].text = shuttle[0].text.replace(target, replace)
                else:
                    replace_begin_index = full_text.index(target)
                    replace_end_index = replace_begin_index + len(target)
                    replace_end_index_in_last_run = replace_end_index - len(''.join(run.text for run in shuttle[:-1]))
                    shuttle[0].text = shuttle[0].text[:replace_begin_index] + replace

                    # Clear middle runs
                    for i in shuttle[1:-1]:
                        i.text = ''

                    # Keep last run
                    shuttle[-1].text = shuttle[-1].text[replace_end_index_in_last_run:]

                # Reset the beginning index for the next iteration
                begin = end + 1
                full_text = ""
    return


def docx_replace(doc, target, replace):
    warning = ''
    try:
        # Replace text in tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        # Replace text in table cells
                        replace_text(paragraph.runs, target, replace)
                        # Replace text in hyperlinks within table cells
                        for link in paragraph._element.xpath(".//w:hyperlink"):
                            replace_text(link.xpath("w:r", namespaces=link.nsmap), target, replace)

        # Replace text in headers and footers
        for section in doc.sections:
            # Process the header and footer for the main body and first page
            for header in [section.header, section.first_page_header]:
                for paragraph in header.paragraphs:
                    # Find all <w:t> tags within the header paragraphs
                    w_t_tags = paragraph._element.xpath(".//w:t")

                    # Replace text in header paragraphs
                    replace_text(w_t_tags, target, replace)

                for table in header.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                # Replace text in table cells within headers
                                replace_text(paragraph.runs, target, replace)
                                # Replace text in hyperlinks within table cells within headers
                                for link in paragraph._element.xpath(".//w:hyperlink"):
                                    replace_text(link.xpath("w:r", namespaces=link.nsmap), target, replace)

            for footer in [section.footer, section.first_page_footer]:
                for paragraph in footer.paragraphs:
                    # Replace text in footer paragraphs
                    replace_text(paragraph.runs, target, replace)

                for table in footer.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                # Replace text in table cells within footers
                                replace_text(paragraph.runs, target, replace)
                                # Replace text in hyperlinks within table cells within footers
                                for link in paragraph._element.xpath(".//w:hyperlink"):
                                    replace_text(link.xpath("w:r", namespaces=link.nsmap), target, replace)

        # Replace text in the main document body
        for paragraph in doc.paragraphs:
            # Replace text in main body paragraphs
            replace_text(paragraph.runs, target, replace)
            # Replace text in hyperlinks within main body paragraphs
            for link in paragraph._element.xpath(".//w:hyperlink"):
                replace_text(link.xpath("w:r", namespaces=link.nsmap), target, replace)
    except IndexError:
        warning += 'File is skipped because of an index error!'

    return warning


"""def text_rebrand(file_in, config):
    # Store file path from CL Arguments.
    file_path = file_in

    if file_path.endswith('.docx'):
        doc = Document(file_path)

        # initialize variable for tables
        tables = doc.tables
        # initialize variable for counting text match occurrences
        occurrences = {}

        for replace_duo in config["ReplaceString"]:
            # initialize the number of occurrences to 0
            occurrences[replace_duo[0]] = 0
            # Loop through sections for header footer content
            for section in doc.sections:
                # Check text of first page header
                fp_header = section.first_page_header

                for fp_header_par in fp_header.paragraphs:
                    # check if paragraph contains text
                    if fp_header_par.text:
                        # check if target text exists in paragraph text
                        if replace_duo[0] in fp_header_par.text:
                            text = fp_header_par.text.replace(replace_duo[0], replace_duo[1])
                            if text != fp_header_par.text:
                                # Replace the text and increment the number of occurrences
                                fp_header_par.text = text
                                occurrences[replace_duo[0]] += 1

                # Check text of first page footer
                fp_footer = section.first_page_footer

                for fp_footer_par in fp_footer.paragraphs:
                    # check if paragraph contains text
                    if fp_footer_par.text:
                        # check if target text exists in paragraph text
                        if replace_duo[0] in fp_footer_par.text:
                            text = fp_footer_par.text.replace(replace_duo[0], replace_duo[1])
                            if text != fp_footer_par.text:
                                # Replace the text and increment the number of occurrences
                                fp_footer_par.text = text
                                occurrences[replace_duo[0]] += 1

        # Loop through tables in document
        for table in tables:
            # Loop through rows in table
            for row in table.rows:
                # Loop through cells in row
                for cell in row.cells:
                    # Loop through paragraphs in cell
                    for paragraph in cell.paragraphs:
                        # check if paragraph contains text
                        if paragraph.text:
                            # check if target text exists in paragraph text
                            if replace_duo[0] in paragraph.text:
                                text = paragraph.text.replace(replace_duo[0], replace_duo[1])
                                # Check if replaced text is not the same as original
                                if text != paragraph.text:
                                    # Replace the text and increment the number of occurrences
                                    paragraph.text = text
                                    occurrences[replace_duo[0]] += 1

        # Loop through paragraphs
        for para in doc.paragraphs:
            # check if paragraph contains text
            if para.text:
                # check if target text exists in paragraph text
                if replace_duo[0] in para.text:
                    text = para.text.replace(replace_duo[0], replace_duo[1])
                    # Check if replaced text is not the same as original
                    if text != para.text:
                        # Replace the text and increment the number of occurrences
                        para.text = text
                        occurrences[replace_duo[0]] += 1

            # Loop through hyperlinks in document
            for link in para._element.xpath(".//w:hyperlink"):
                inner_run = link.xpath("w:r", namespaces=link.nsmap)[0]

                # Check if hyperlink contains text
                if inner_run.text:
                    # Check if hyperlink text contains target tex
                    if replace_duo[0] in inner_run.text:
                        text = inner_run.text.replace(replace_duo[0], replace_duo[1])
                        # Check if replaced text is not the same as original
                        if text != inner_run.text:
                            # Replace the text and increment the number of occurrences
                            inner_run.text = text
                            occurrences[replace_duo[0]] += 1

    # print the number of occurrences of each word
    for word, count in occurrences.items():
        if count > 0:
            print(f"{word} ({count})")

    # save the new docx file
    new_file_path = os.path.basename(file_path)
    doc.save((config["OutputFolder"]) + new_file_path)

    return"""

"""def paragraph_replace_text(paragraph, regex, replace_str):
    ""
    Return `paragraph` after replacing all matches for `regex` with `replace_str`.

    `regex` is a compiled regular expression prepared with `re.compile(pattern)`
    according to the Python library documentation for the `re` module.
    ""
    # --- a paragraph may contain more than one match, loop until all are replaced ---
    while True:
        text = paragraph.text
        match = regex.search(text)
        if not match:
            break

        # --- when there's a match, we need to modify run.text for each run that
        # --- contains any part of the match-string.
        runs = iter(paragraph.runs)
        start, end = match.start(), match.end()

        # --- Skip over any leading runs that do not contain the match ---
        for run in runs:
            run_len = len(run.text)
            if start < run_len:
                break
            start, end = start - run_len, end - run_len

        # --- Match starts somewhere in the current run. Replace match-str prefix
        # --- occurring in this run with entire replacement str.
        run_text = run.text
        run_len = len(run_text)
        run.text = "%s%s%s" % (run_text[:start], replace_str, run_text[end:])
        end -= run_len  # --- note this is run-len before replacement ---

        # --- Remove any suffix of match word that occurs in following runs. Note that
        # --- such a suffix will always begin at the first character of the run. Also
        # --- note a suffix can span one or more entire following runs.
        for run in runs:  # --- next and remaining runs, uses same iterator ---
            if end <= 0:
                break
            run_text = run.text
            run_len = len(run_text)
            run.text = run_text[end:]
            end -= run_len

    # --- optionally get rid of any "spanned" runs that are now empty. This
    # --- could potentially delete things like inline pictures, so use your judgement.
    # for run in paragraph.runs:
    #     if run.text == "":
    #         r = run._r
    #         r.getparent().remove(r)

    return paragraph"""


def copy_and_replace(zip_in, zip_out):
    """
    Copies the zip file except for the media folder

    Parameters
    ----------
    zip_in: ZipFile obj
        input ZipFile

    zip_out: ZipFile obj
        output ZipFile
    """
    # XML files that contain image crop data
    xml_name_crop = ['header1.xml', 'header2.xml', 'header3.xml']
    # Go over every file in the input document
    for path in zip_in.namelist():
        already_added = False
        # Check if the path is not in the media folder
        if "media" not in path:
            file_content = zip_in.read(path)
            # Remove crop from images
            if any(element in path for element in xml_name_crop):
                already_added = _modify_xml_image_crop_fit(zip_in, zip_out, path)

            # Copy file into new document
            if not already_added:
                zip_out.writestr(path, file_content)


"""def try_decode(content):
    encodings = ['utf-8', 'latin-1']  # Add more encodings if necessary
    for encoding in encodings:
        try:
            decoded_content = content.decode(encoding)
            return decoded_content
        except UnicodeDecodeError:
            continue
    return None"""


def resize_image(input_image_path, output_image_folder, reference_image_path):
    """
    Resize the input image to match the aspect ratio of the reference image, preserving the aspect ratio and avoiding
    stretching.

    Parameters:
        input_image_path (str): Path to the input image file.
        output_image_folder (str): Folder path to save the resized image.
        reference_image_path (str): Path to the reference image for aspect ratio.

    Returns:
        new_image_path (str): Path to the new resized image.
    """
    # Catch any unsupported images
    try:
        # Open the reference image to get its aspect ratio
        with Image.open(reference_image_path) as reference_image:
            ref_height = reference_image.height
            ref_width = reference_image.width
            ref_aspect_ratio = ref_width / ref_height

        new_image_path = f'{output_image_folder}/resized_image_{ref_width}x{ref_height}.png'
        if not os.path.exists(new_image_path):

            # Open the input image
            with Image.open(input_image_path) as input_image:
                # Calculate the new dimensions for the input image while maintaining aspect ratio
                input_aspect_ratio = input_image.width / input_image.height

                if input_aspect_ratio > ref_aspect_ratio:
                    # The input image is wider, adjust the width to match the reference aspect ratio
                    new_width = ref_width
                    new_height = int(new_width / input_aspect_ratio)
                else:
                    # The input image is taller, adjust the height to match the reference aspect ratio
                    new_height = ref_height
                    new_width = int(new_height * input_aspect_ratio)

                # Resize the input image with aspect ratio preserved
                resized_image = input_image.resize((new_width, new_height))
                # Create a blank image with the reference aspect ratio and paste the resized image onto it
                result_image = Image.new("RGB", (ref_width, ref_height), (255, 255, 255))
                offset = ((ref_width - new_width) // 2, (ref_height - new_height) // 2)

                result_image.paste(resized_image, offset)

                # Save the resized image
                result_image.save(new_image_path)
    except Exception:
        new_image_path = ''

    return new_image_path


def replace_header_images(zip_in, zip_out, config, note):
    """
    Search for images in the header and replace them with the logo

    Parameters
    ----------
    zip_in: ZipFile obj
        input ZipFile

    zip_out: ZipFile obj
        output ZipFile

    config: dictionary
        dict containing information from the config file

    note: string
        notes about the location of the logo

    Returns
    -------
    note: string
        notes about the location of the logo

    warning: string
        warning if multiple header images have been replaced

    images_replaced_path: string array
        paths to the images replaced
    """
    warning = ""
    images_replaced_path = []
    # Check every filename for header --> can contain information about images in the header
    # Possible filenames: header2.xml.rels, header3.xml.rels
    for file in zip_in.namelist():
        if "header" in file and "rel" in file:
            # pylint: disable=W1401
            # Disable error from using \S\s in a binary string which is needed for regex
            # Get every target image
            image_locations = findall(b'Target="media[\S\s]*?"', zip_in.read(file))

            if len(image_locations) > 1:
                warning = "Warning: Multiple images found in header"
            for image_location in image_locations:
                image_location = image_location[image_location.find(b"media"):-1].decode("ascii")
                # Get the full path: Path in target string is just media/[image]
                # Path is either word/media/[image] or xls/media/[image]
                try:
                    # Skip image if not in the right format
                    if image_location.endswith('.wmf'):
                        continue

                    header_image_path = f"{config['filetype']}/{image_location}"

                    # Add replaced image path for output
                    images_replaced_path.append(header_image_path)

                    # Extract image from the zip object
                    zip_image_location = 'word/' + image_location
                    zip_in.extract(zip_image_location, path=config['ImagesFolder'])

                    # Resize image to fit in the extracted image's container
                    new_image_path = resize_image(config['NewLogoPath'], config['ImagesFolder'],
                                                  config['ImagesFolder'] + zip_image_location)

                    # Delete extracted image
                    os.remove(config['ImagesFolder'] + zip_image_location)

                    if header_image_path not in zip_out.namelist():
                        zip_out.write(new_image_path, header_image_path)
                        log_information['LegacyText'] += f"Replaced header image {header_image_path},"
                except KeyError:
                    # Filetype is not in config
                    # Get path to the image by looping over the entire content
                    for location in zip_in.namelist():
                        if image_location in location:
                            if image_location not in zip_out.namelist():
                                print(f"Replaced header image at {location}")
                                zip_out.write(new_image_path, location)
                                log_information['LegacyText'] += f"Replaced header image {location},"
                            break

    return note, warning, images_replaced_path


def check_header_images(zip_in, config, note):
    """
    checks header images and adds information to note

    Parameters
    ----------
    zip_in: ZipFile obj
        input ZipFile

    config: dictionary
        dict containing information from the config file

    note: String
        Information about the logo location

    Returns
    -------
    note: String
        Information about the logo location and possible alternatives
    """
    alternative_old_logo_found = False
    # Found image at expected path but with different size
    # Check if image is in header
    for file in zip_in.namelist():
        if "header" in file and "rel" in file:
            # pylint: disable=W1401
            # Disable error from using \S\s in a binary string which is needed for regex
            # Get every target image
            image_locations = findall(b'Target="media[\S\s]*?"', zip_in.read(file))
            for image_location in image_locations:
                image_location = image_location[image_location.find(
                    b"media"):-1].decode("ascii")
                # Format image_location to the full path
                try:
                    image_location = f'{config["filetype"]}/{image_location}'
                except KeyError:
                    # filetype not in config; cant automatically create full image path
                    # Check every location to get the full path
                    for path in zip_in.namelist():
                        if image_location in path:
                            image_location = path
                            break

                alternative_old_logo_size = zip_in.getinfo(
                    image_location).file_size
                note += f'''Found alternative image in {image_location
                } with size {alternative_old_logo_size} bytes, '''
                alternative_old_logo_found = True
    if not alternative_old_logo_found:
        note += "No image in header,"
    return note


def place_logo_header(zip_in, zip_out, config):
    """
    Places the new BH logo in the Word document

    Parameters
    ----------
    zip_in: ZipFile obj
        input ZipFile

    zip_out: ZipFile obj
        output ZipFile

    config: dictionary
        dict containing information from the config file

    Returns
    -------
    status: string
        status information about the finding of the logo

    note: string
        notes about the location of the logo

    warning: string
        warning if multiple header images have been replaced
    """

    status = "LogoNotFound"
    note = ""
    warning = ""
    headercount = 0
    header_image_paths = []

    log_information['NumLogos'] = 0

    # Check if every image in the header should be replaced
    if "true" in config["ReplaceHeaderImage"].lower():
        note, warning, header_image_paths = replace_header_images(zip_in, zip_out, config, note)
    else:
        note = check_header_images(zip_in, config, note)

    if len(header_image_paths) == 0:
        log_information['Logo'] = "No Header Logo"
    else:
        log_information['Logo'] = "Header Logo Found"

    file_name = os.path.basename(zip_out.filename)
    header_images[file_name] = header_image_paths
    # Filter the file paths that are in 'word/media' 
    files_in_folder = [name for name in zip_in.namelist() if name.startswith('word/media')]

    # Get the number of files in the 'word/media' folder
    num_files = len(files_in_folder)

    # Add header image(s) to the catalog (if they're unique)
    if config['CompareLogoByPixels'] and num_files > 0:
        for image_location in header_images[file_name]:
            # Add image to the catalog
            add_image_to_catalog(zip_in, image_location, config)

    return status, note, warning


def add_image_to_catalog(zip_in, image_location, config):
    """
    Extracts the image from a zipfile path, and adds it to the catalog of logos if it's unique.

    Parameters
    ----------
    zip_in: ZipFile obj
        input ZipFile

    image_location: string
        path to the image inside the zipfile

    config: dictionary
        dict containing information from the config file
    """
    global different_logos_found
    # Extract image
    zip_in.extract(image_location, path=config["BetweenFolder"])
    file_extension = os.path.splitext(config["BetweenFolder"] + image_location)[1]
    logo_is_present = False
    image_is_transparent = False

    # Check if image is transparent and omit it if it is
    try:
        image = Image.open(config["BetweenFolder"] + image_location)
        if image.mode == 'RGBA':
            image_is_transparent = all(pixel[3] == 0 for pixel in image.getdata())
        try:
            image.close()
        except Exception as e:
            print(f'Error when closing the image. Error: {e}')
    except Exception as e:
        print(f'Error occurred when checking image transparency. Error: {e}')

    if not image_is_transparent:
        # Get all logo paths from the catalog
        logo_locations = os.listdir(config["FoundLogosFolder"])

        # Cycle through the entire catalog and check if header image is already present
        for logo in logo_locations:
            # Compare header image with logo
            logo_is_present = compare_images(config["BetweenFolder"] + image_location,
                                             config["FoundLogosFolder"] + logo)
            if logo_is_present:
                break

        # Add image to the logo catalog folder
        if not logo_is_present:
            different_logos_found += 1
            renamed_path = f'{config["BetweenFolder"]}{os.path.dirname(image_location)}/logo_' \
                           f'{different_logos_found}{file_extension}'
            os.rename(config["BetweenFolder"] + image_location, renamed_path)
            shutil.move(renamed_path, config["FoundLogosFolder"])

        # Remove zip images empty base directory
        zip_image_base_directory = image_location.split('/')[0]
        shutil.rmtree(config["BetweenFolder"] + zip_image_base_directory)


def place_logo_body(file_in, file_out, config):
    log_information['Inputfile'] = file_in
    file_path = file_in
    header_images_paths = {}

    new_file_path = os.path.basename(file_path)
    skip_file_formats = ('.pptx', '.pdf', '.xlsx', '.xls')

    if file_in.endswith('.docx'):
        doc = Document(file_path)
        header_images_paths = header_images[new_file_path]
        doc.save((config["BetweenFolder"]) + new_file_path)
    elif file_in.endswith(skip_file_formats):
        # Move file and finish executing function
        os.rename(file_in, file_out)

        return True

    # Open input document and new document
    with ZipFile(open((config["BetweenFolder"]) + new_file_path, "rb")) as zip_in:
        with ZipFile(file_out, "w", ZIP_DEFLATED) as zip_out:

            # Copy contents to zip_out
            for path in zip_in.namelist():
                if 'media' not in path:
                    file_content = zip_in.read(path)
                    # Copy file into new document
                    zip_out.writestr(path, file_content)

            # Get all image locations without the header
            image_locations = [item for item in zip_in.namelist() if '/media/' in item]

            # Check that header images exist
            if len(header_images_paths) != 0:
                # Remove from image_location header paths
                for header_image in header_images_paths:
                    if header_image in image_locations:
                        image_locations.remove(header_image)

            # Extract all images besides the header logo
            zip_in.extractall(config["BetweenFolder"], members=image_locations)

            # Get all logo paths from the catalog
            logo_locations = os.listdir(config["FoundLogosFolder"])

            # Check that there's any images
            if len(image_locations) != 0:
                # Compare all images with the logo catalog
                for image in image_locations:
                    for logo in logo_locations:
                        similar = compare_images(config["FoundLogosFolder"] + logo, config["BetweenFolder"] + image)

                        # Replace image if similar
                        if similar:
                            log_information['Logo'] = 'Body Logo Found'
                            log_information['NumLogos'] += 1
                            zip_image_location = f'{os.path.dirname(image)}/{os.path.basename(image)}'

                            # Resize logo from catalog
                            logo_replacement_path = resize_image(config['NewLogoPath'], config['ImagesFolder'],
                                                                 config["BetweenFolder"] + image)

                            # Add similar images to zip_out
                            zip_out.write(logo_replacement_path, zip_image_location, compress_type=ZIP_DEFLATED)

                            # Add note
                            log_information['LegacyText'] += f'Replaced image: {image} '

                            break

            if len(image_locations) != 0:
                zip_image_base_directory = image.split('/')[0]
            elif len(header_images_paths) != 0:
                zip_image_base_directory = header_images_paths[0].split('/')[0]
            else:
                zip_image_base_directory = None

            # Delete extracted images
            if zip_image_base_directory is not None:
                if os.path.exists(config["BetweenFolder"] + zip_image_base_directory):
                    shutil.rmtree(config["BetweenFolder"] + zip_image_base_directory)

            # Add missing images
            add_missing_images(zip_in, zip_out)
    # Remove file from betweenFolder
    os.remove((config["BetweenFolder"]) + new_file_path)
    # Remove file from original location
    os.remove(file_in)

    log_information['Notes'] = 'File processed successfully - body image'
    # log.write(f'{file_in};{logo_found};{str(num_logos)};{status};{note};{warning}\n')


def get_filetype(file, config):
    """
    Returns the config containing information about the current filetype.
    When it's neither a Word file nor an Excel file the config returned
    won't contain information about the filetype

    Parameters
    ----------
    file: string
        name of the file to be checked

    config: dictionary
        content of the configuration file as dictionary with the tag as
        key and the information as value

    Returns
    -------
    config: dictionary
        content of the configuration file as dictionary with the tag as
        key and the information as value with added filetype information
    """
    if ".do" in file:
        config["filetype"] = "word"
    elif ".xls" in file:
        config["filetype"] = "xl"
    elif ".ppt" in file:
        config["filetype"] = "ppt"
    elif ".pdf" in file:
        config["filetype"] = "pdf"
    else:
        # Remove "filetype" from config if filetype is unknown
        try:
            config.pop("filetype")
        except KeyError:
            pass
    return config


def get_filetypes_in_folder(folder):
    """
    Returns a list of all filetypes (Word, Excel or PowerPoint) in a folder

    Parameters
    ----------
    folder: string
        the path to the folder containing the files

    Returns
    -------
    filetype: list
        list of the filetypes (Word,Excel or PowerPoint) found in the folder
    """
    filetypes = []
    for file in os.listdir(folder):
        if ".do" in file:
            if "Word" not in filetypes:
                filetypes.append("Word")
        elif ".xls" in file:
            if "Excel" not in filetypes:
                filetypes.append("Excel")
        elif ".ppt" in file:
            if "PowerPoint" not in filetypes:
                filetypes.append("PowerPoint")
    return filetypes


def start_com_servers(filetypes):
    """
    Starts the needed COM servers for PDF conversion

    Parameters
    ----------
    filetypes: list
        list of filetypes of the files which are going to be converted
    """
    word = None
    excel = None
    power_point = None
    if "Word" in filetypes:
        print("Starting Word COM server")
        word = comtypes.client.CreateObject("Word.Application")
        word.Visible = False
    if "Excel" in filetypes:
        print("Starting Excel COM server")
        excel = comtypes.client.CreateObject("Excel.Application")
        excel.Visible = False
    if "PowerPoint" in filetypes:
        print("Starting PowerPoint COM server")
        power_point = comtypes.client.CreateObject("PowerPoint.Application")
        # Trying to hide the window will result in the program crashing
        # power_point.Visible = False
    return word, excel, power_point


def quit_com_servers(word, excel, power_point):
    """
    Stops all active COM servers

    Parameters
    ----------
    word: comtypes.Pointer
        pointer to the Word COM server

    excel: comtypes Pointer to the Excel COM server

    power_point: comtypes.Pointer to the PowerPoint COM server
    """
    print("Stopping COM server")
    if word is not None:
        word.Quit()
    if excel is not None:
        excel.Quit()
    if power_point is not None:
        power_point.Quit()


def convert_to_pdf(file_in, config, word, excel, power_point):
    """
    converts a file into a pdf file and saves it

    Parameters
    ----------
    file_in: string
        path to the input file

    config: dict
        contains information set in the config file and the current filetype

    word: comtypes.Pointer to the Word COM server

    excel: comtypes.Pointer to the Excel COM server

    power_point: comtypes.Pointer to the PowerPoint COM server
    """
    # Use the correct COM server
    if config["filetype"] == "word":
        doc = word.Documents.Open(file_in)
        doc.SaveAs(config["PDFPath"], FileFormat=FILE_FORMAT_PDF_WORD)
        doc.Close()
    elif config["filetype"] == "xl":
        workbook = excel.Workbooks.Open(file_in)
        workbook.ExportAsFixedFormat(FILE_FORMAT_PDF_EXCEL, config["PDFPath"])
        workbook.close()
    elif config["filetype"] == "ppt":
        presentation = power_point.Presentations.Open(file_in)
        presentation.ExportAsFixedFormat(config["PDFPath"], FILE_FORMAT_PDF_PPT)
        presentation.close()


def process_file_word(file_in, file_out, config):
    # log_information['Inputfile'] = file_in
    file_path = file_in
    file_out_path = file_out

    # log variables
    # status, note, warning, text_note = '', '', '', ''

    # Convert file to docx if it ends in doc
    if file_path.endswith('.doc') or file_path.endswith('.docm'):
        # Convert .doc file to .docx format and update file_path
        file_path = convert_file(file_path, FILE_FORMAT_DOCX)

        # Update file_out_path to contain the new extension
        file_out_path = f'{os.path.dirname(file_out)}\{os.path.basename(file_path)}'

    doc = Document(file_path)
    prop = doc.core_properties
    for replace_duo in config["ReplaceString"]:
        log_information['Warning'] = docx_replace(doc, replace_duo[0], replace_duo[1])
        prop.title = prop.title.replace(replace_duo[0], replace_duo[1])

    new_file_path = os.path.basename(file_path)
    doc.save((config["BetweenFolder"]) + new_file_path)

    # Open input document and new document
    with ZipFile(open((config["BetweenFolder"]) + new_file_path, "rb")) as zip_in:
        with ZipFile(file_out_path, "w", ZIP_DEFLATED) as zip_out:
            # copy document and replace content
            copy_and_replace(zip_in, zip_out)
            # check for logo
            status, note, warning = place_logo_header(zip_in, zip_out, config)
            # Add missing images
            add_missing_images(zip_in, zip_out)

    # Remove file from betweenFolder
    os.remove((config["BetweenFolder"]) + new_file_path)

    log_information['Notes'] = 'File processed successfully - text & header'

    # log.write(f"{file_in};{status};-;{note};{text_note};{warning}\n")


def copy_and_replace_content_excel(file_in, file_out, config):
    """
    Perform copying the content from the input file, replaces imagery and text

    Parameters
    ----------
    file_in: string
        path to the input file
    file_out: string
        path to the output file
    config: dict
        Configuration file

    Returns None
    -------
    """
    new_file_path = os.path.basename(file_in)

    # Open the input Zip file for reading and the output Zip file for writing
    with ZipFile(open(config['InputFolder'] + '\\' + new_file_path, 'rb')) as zip_in:
        with ZipFile(file_out, 'w') as zip_out:

            # Copy the content of the input file except the media folder and xml files containing strings to be replaced
            # on the worksheets
            for path in zip_in.namelist():
                if 'media' not in path:
                    if 'worksheets' in path or 'sharedStrings' in path:
                        if path.endswith('xml'):
                            continue
                    file_content = zip_in.read(path)
                    zip_out.writestr(path, file_content)

            # Extract all images from the workbook stored in the media folder
            image_locations = [img_loc for img_loc in zip_in.namelist() if '/media/' in img_loc]
            # Extract all images
            zip_in.extractall(config['BetweenFolder'], members=image_locations)

            # Get all logo paths from the catalog
            logo_locations = os.listdir(config['FoundLogosFolder'])

            # Check if there are any images in the presentation
            if len(image_locations) > 0:
                # Compare all images with the logo catalog
                for image_location in image_locations:
                    for logo_location in logo_locations:
                        # Do similarity check
                        similarity = compare_images(
                            image_path1=config['FoundLogosFolder'] + logo_location,  # Path for the first image
                            image_path2=config['BetweenFolder'] + image_location  # Path for the second image
                        )

                        # If images are similar, replace them
                        if similarity:
                            log_information['Logo'] = 'Logo Found'
                            log_information['NumLogos'] += 1
                            zip_image_location = f'{os.path.dirname(image_location)}/{os.path.basename(image_location)}'

                            # Resize the logo from catalog
                            resized_image_path = resize_image(
                                input_image_path=config['NewLogoPath'],  # Path to 'replacementLogo.png'
                                output_image_folder=config['ImagesFolder'],  # Path to save resized image
                                reference_image_path=config['BetweenFolder'] + image_location  # Path to reference image
                            )

                            # Add new resized image to archive under the original image file name
                            zip_out.write(resized_image_path, zip_image_location, compress_type=ZIP_DEFLATED)

                            log_information['LegacyText'] += f'Replaced image: {image_location}, '

                            break  # Break out of the inner loop

                    else:  # If no similarity detected add the image file to the output archive
                        if image_location not in zip_out.namelist():
                            image_data = zip_in.read(image_location)
                            zip_out.writestr(image_location, image_data)

            # Perform text replacement in the xml files
            _replace_text_excel(zip_in=zip_in, zip_out=zip_out, config=config)

            log_information['Notes'] = 'File processed successfully'
            # log.write(f'{file_in};{logo};{str(num_logos)};{status};{note};{warning}\n')


def _replace_text_excel(zip_in, zip_out, config):
    """
    Finds elements in the input files that contain outdated text and replaces them

    Parameters
    ----------
    zip_in: Input Zip archive
    zip_out: Output Zip archive
    config: dict - Configuration file

    Returns None
    -------
    """

    prefix_map = {
        'xmlns': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'
    }  # Mapping of namespaces in the XML

    # Iterate through all files in the input archive
    for path in zip_in.namelist():
        # If worksheets or the sharedStrings XML files are present in the archive
        if '/worksheets/' in path or 'sharedStrings' in path:
            if path.endswith('xml'):

                xml_tree = zip_in.read(path)  # Open and read in the XML file from path
                root = lxml.etree.fromstring(xml_tree)  # Parse the XML tree

                # Variable to hold xml elements that will need to be string validated
                string_elements = []

                if 'sharedStrings' in path:
                    # Find all t - 'text' tags in the sharedStrings XML file
                    string_elements = root.findall('.//xmlns:t', namespaces=prefix_map)
                else:
                    # In the worksheet XML find 'oddHeader' and 'oddFooter' elements that might contain strings that
                    # need to be replaced
                    odd_headers = root.findall('.//xmlns:oddHeader', namespaces=prefix_map)
                    string_elements.extend(odd_headers)
                    odd_footers = root.findall('.//xmlns:oddFooter', namespaces=prefix_map)
                    string_elements.extend(odd_footers)

                # If elements with text were found pass them to the replace string function
                if string_elements is not None:
                    _replace_string_excel(elements=string_elements, config=config)

                # Write the XML with replaced text to the output file at the path
                replaced_text_xml = lxml.etree.tostring(root, encoding='UTF-8')
                zip_out.writestr(path, replaced_text_xml)


def _replace_string_excel(elements, config):
    """
    Performs text replacement on XML elements with text values

    Parameters
    ----------
    elements: List
        List of XML elements with text value to be replaced
    config: dict
        Configuration file

    Returns None
    -------
    """

    for element in elements:
        for value in config['ReplaceString']:
            replaced_value = element.text.replace(value[0], value[1])
            if replaced_value != element.text:
                element.text = replaced_value


def process_file_excel(file_in, file_out, config):
    file_input_path = file_in  # Input file path

    # If file in wrong format then convert .xls or .xlsm file to .xlsx and update file path
    if file_input_path.endswith('.xls'):  # or file_input_path.endswith('.xlsm'):
        file_input_path = convert_file(file_in, FILE_FORMAT_XLSX)
        # Update the output basename accordingly, so it will be saved as 'xlsx'
        file_out = file_out.replace(file_out.rsplit(".")[-1], 'xlsx')

    copy_and_replace_content_excel(file_in=file_input_path, file_out=file_out, config=config)

    return True


def pptx_replace(prs, config):
    """
    Replace all text in the presentation slides and notes section

    Parameters
    ----------
    prs: The PowerPoint Presentation object
    config: Configuration file

    Returns None
    -------

    """

    for slide in prs.slides:  # Iterate over each slide in the presentation
        for shape in slide.shapes:  # Iterate through the building blocks (shapes) of the slide
            try:
                text_frame = shape.text_frame  # Get the text frame from the slide
                # Call the replace function passing in the frame and the configuration file
                _pptx_text_replace(text_frame=text_frame, config=config)
            except AttributeError:  # Certain objects have no text frame attribute hence they will be skipped
                continue
        # Check whether the slide contains notes
        if slide.has_notes_slide:
            # If the notes section is not empty get the text frame
            if slide.notes_slide.notes_text_frame:
                notes_text_frame = slide.notes_slide.notes_text_frame  # Get the text frame from the notes slide
                # Call the replace function passing in the frame and the configuration file
                _pptx_text_replace(text_frame=notes_text_frame, config=config)


def _pptx_text_replace(text_frame, config):
    """
    Calls the 'replace_text' function an applies text replacement for the PowerPoint text frame

    Parameters
    ----------
    text_frame: A PowerPoint text frame either from a slide or from a note
    config: Configuration file containing string replacement definitions

    Returns None
    ------
    """

    # Iterate through each replacement definitions in the config file
    for replace_duo in config['ReplaceString']:
        # Iterate through each 'paragraph' block in the text frame
        for paragraph in text_frame.paragraphs:
            # Call the 'replace_text' function passing in the 'runs' from the paragraph, the target and replacement
            # strings
            replace_text(runs=paragraph.runs, target=replace_duo[0], replace=replace_duo[1])


def process_file_powerpoint(file_in, file_out, config):
    """
    Replaces old Logo / text in PowerPoint files

    Parameters
    ----------
    file_in: string
        path to the input file
    file_out: string
        path to the output file
    config: dict
        Key-value pairs for configuration in a form of a dictionary

    Returns None
    -------
    """
    # log_information['Inputfile'] = file_in

    file_path = file_in  # Path to the input file
    file_out_path = file_out  # Path to the output file

    prs = Presentation(file_path)  # Instantiate the Presentation object using the python-pptx module

    pptx_replace(prs=prs, config=config)  # Replace text in presentation slides and notes

    new_file_path = os.path.basename(file_path)
    prs.save((config['BetweenFolder']) + new_file_path)

    # Open the input document to access the underlying XML files for the PowerPoint presentation
    with ZipFile(open(config['BetweenFolder'] + new_file_path, 'rb')) as zip_in:
        with ZipFile(file_out_path, 'w', ZIP_DEFLATED) as zip_out:

            # Copy the content of the input file to the output file except the media folder and slide xml files
            for path in zip_in.namelist():
                if 'media' not in path:
                    if '/slides/' in path:
                        if path.endswith('.xml'):
                            continue
                    file_content = zip_in.read(path)
                    zip_out.writestr(path, file_content)

            # Extract all images from the presentation stored in the media folder
            image_locations = [img_loc for img_loc in zip_in.namelist() if '/media/' in img_loc]
            # Extract all images
            zip_in.extractall(config['BetweenFolder'], members=image_locations)

            # Get all logo paths from the catalog
            logo_locations = os.listdir(config['FoundLogosFolder'])

            # Check if there are any images in the presentation
            if len(image_locations) > 0:
                # Compare all images with the logo catalog
                for image_location in image_locations:
                    # print(image_location)
                    for logo_location in logo_locations:
                        # Do similarity check
                        similarity = compare_images(
                            image_path1=config['FoundLogosFolder'] + logo_location,  # Path for the first image
                            image_path2=config['BetweenFolder'] + image_location  # Path for the second image
                        )

                        # If images are similar, replace them
                        if similarity:
                            log_information['Logo'] = 'Logo Found'
                            log_information['NumLogos'] += 1
                            zip_image_location = f'{os.path.dirname(image_location)}/{os.path.basename(image_location)}'

                            # Resize the logo from catalog
                            resized_image_path = resize_image(
                                input_image_path=config['NewLogoPath'],  # Path to 'replacementLogo.png'
                                output_image_folder=config['ImagesFolder'],  # Path to save resized image
                                reference_image_path=config['BetweenFolder'] + image_location  # Path to reference image
                            )

                            # Add new resized image to archive under the original image file name
                            zip_out.write(resized_image_path, zip_image_location, compress_type=ZIP_DEFLATED)

                            log_information['LegacyText'] += f'Replaced image: {image_location}, '

                            break  # Break out of the inner loop

                    else:  # If no similarity detected add the image file to the output archive
                        if image_location not in zip_out.namelist():
                            image_data = zip_in.read(image_location)
                            zip_out.writestr(image_location, image_data)

                # Get base directory where images are located
                zip_image_base_directory = image_location.split('/')[0]

                # Delete extracted images following iteration
                if os.path.exists(config['BetweenFolder'] + zip_image_base_directory):
                    shutil.rmtree(config['BetweenFolder'] + zip_image_base_directory)

            for path in zip_in.namelist():
                if '/slides/' in path and path.endswith('.xml'):
                    # For each slide in the PowerPoint file hide the background graphics
                    disable_background_graphics(zip_in=zip_in, zip_out=zip_out, xml_path=path)

    os.remove(config['BetweenFolder'] + new_file_path)  # Remove the file from the 'between' folder

    prs = Presentation(file_out_path)  # Load in the presentation from the output file

    insert_replacement_image_to_slide(prs=prs, config=config)  # Insert the BakerHughes logo in

    change_bg_color(prs=prs)  # Change the background color of the slides in the presentation

    prs.save(file_out_path)  # Save the presentation

    log_information['Notes'] = 'File processed successfully'

    return True


def _check_background_color(slide_layout):
    """
    Checks whether the old BH RGB color scheme (HEX values --> blue) is present on the slide

    Parameters
    ----------
    slide_layout: SlideLayout object used by the current slide

    Returns XML element containing old RGB colors to replace or None
    -------
    """

    prefix_map = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
    }  # Contains the mapping from namespace mapping to full name for searching the XML tree

    old_colors_to_replace = ['005EB8', '00B5E2']  # HEX values of colors to look for on the slide

    if slide_layout:  # If the slide layout exists
        # Find all the elements that contain RGB color properties on the slide layout
        color_lists = slide_layout.background.element.findall('.//a:gsLst', namespaces=prefix_map)
        if color_lists is not None:
            # Iterate over the found elements
            for color_list in color_lists:
                # If the 'Background Properties' element is the parent element
                if 'bgPr' in color_list.getparent().getparent().tag:
                    colors = color_list.findall('.//a:srgbClr', namespaces=prefix_map)  # Extract RGB colors
                    if colors is not None:
                        if colors[0].attrib is not None and colors[1].attrib is not None:
                            # If color attributes are present and equal to the colors to be replaced
                            if colors[0].attrib['val'] in old_colors_to_replace and \
                                    colors[1].attrib['val'] in old_colors_to_replace:
                                return color_list  # Return the 'color list' element
    return None  # Return None otherwise


def change_bg_color(prs):
    """
    Checks whether the old BH RGB color scheme (HEX values --> blue) is present on the slide

    Parameters
    ----------
    prs: Presentation object

    Returns None
    -------
    """

    # NEW COLORS - 05322b --> green, 1d2920 --> darker green, 018374 -> very light green
    # OLD COLORS - 005EB8 & 00B5E2

    color_changed = False

    prefix_map = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
    }  # Contains the mapping from namespace mapping to full name for searching the XML tree

    # Iterate over all slide layouts in the presentation
    for slide_layout in prs.slide_layouts:
        color_list = _check_background_color(slide_layout=slide_layout)  # Check if old BH color scheme is present
        if color_list is not None:
            # Find the elements containing the RGB colors
            colors = color_list.findall('.//a:srgbClr', prefix_map)
            if colors is not None:
                for clr in colors:
                    if clr.attrib is not None:
                        # Change old colors to new BH colors
                        if clr.attrib['val'] == '005EB8':
                            clr.attrib['val'] = '018374'
                        if clr.attrib['val'] == '00B5E2':
                            clr.attrib['val'] = '05322b'

                        color_changed = True

    return color_changed


def _create_copyright_element():
    """
    Defines an XML string then parses it to create an element containing copyright information for BH

    Parameters
    ----------

    Returns XML element
    -------
    """

    # String representation of the XML element
    xml_string = """<p:sp>
                        <p:nvSpPr>
                            <p:cNvPr id="999" name="Copyright Text"/>
                            <p:cNvSpPr txBox="1"/>
                            <p:nvPr userDrawn="1"/>
                        </p:nvSpPr>
                        <p:spPr>
                            <a:xfrm>
                                <a:off x="404874" y="6608933"/>
                                <a:ext cx="2884874" cy="369332"/>
                            </a:xfrm>
                            <a:prstGeom prst="rect">
                                <a:avLst/>
                            </a:prstGeom>
                            <a:noFill/>
                        </p:spPr>
                        <p:txBody>
                            <a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0">
                                <a:spAutoFit/>
                            </a:bodyPr>
                            <a:lstStyle/>
                            <a:p>
                                <a:r>
                                    <a:rPr lang="en-US" sz="800" spc="-10" baseline="0" dirty="0">
                                        <a:solidFill>
                                            <a:schemeClr val="accent4"/>
                                        </a:solidFill>
                                    </a:rPr>
                                    <a:t>Copyright 2021 Baker Hughes Company. All rights reserved.</a:t>
                                </a:r>
                            </a:p>
                        </p:txBody>
                    </p:sp>"""

    parser = lxml.etree.XMLParser(encoding='UTF-8', recover=True, remove_blank_text=True)  # XML parser object
    # Parse the string definition and create the XML element
    copyright_xml_element = xml.etree.ElementTree.fromstring(xml_string, parser=parser)

    return copyright_xml_element  # Return the XML element


def insert_replacement_image_to_slide(prs, config):
    """
    Inserts the replacement BH logo to each slide in the PowerPoint presentation

    Parameters
    ----------
    prs: Presentation object
    config: Configuration file

    Returns None
    -------
    """

    slide_width = prs.slide_width.inches  # Get the slide width in inches
    slide_height = prs.slide_height.inches  # Get the slide height in inches

    # Define possible logo positions (top-right, bottom-left)
    logo_positions = {
        'bottom-right': {
            'left': slide_width - 2.55,  # Distance from left edge of slide
            'top': slide_height - 0.65  # Distance from the top of the slide
        },
        'top-left-first': {
            'left': 0.6,
            'top': 0.3
        },
        'bottom-left': {
            'left': 0.5,
            'top': slide_height - 0.6
        }
    }

    # Iterate through all slides in the Presentation
    for slide_num, slide in enumerate(prs.slides):
        if slide_num == 0:  # If the slide is the first slide (Title slide)
            # Check if old color scheme exists on the current slide's slide layout
            if _check_background_color(slide_layout=slide.slide_layout) is not None:
                # Add the Large, White BH logo to the top left of the slide
                slide.shapes.add_picture(image_file=config['ReplacementImageWhiteLarge'],
                                         left=pptx.util.Inches(logo_positions['top-left-first']['left']),
                                         top=pptx.util.Inches(logo_positions['top-left-first']['top'])
                                         )
                if check_for_logo_collision_updated(slide):
                    _delete_replacement_image(slide)
                    slide.shapes.add_picture(image_file=config['ReplacementImageWhiteSmall'],
                                             left=pptx.util.Inches(logo_positions['bottom-right']['left']),
                                             top=pptx.util.Inches(logo_positions['bottom-right']['top'])
                                             )
            else:  # If old color scheme is not detected
                # Add the Large, Dark BH logo to the top left of the slide
                slide.shapes.add_picture(image_file=config['ReplacementImageDarkLarge'],
                                         left=pptx.util.Inches(logo_positions['top-left-first']['left']),
                                         top=pptx.util.Inches(logo_positions['top-left-first']['top'])
                                         )
                if check_for_logo_collision_updated(slide):
                    _delete_replacement_image(slide)
                    slide.shapes.add_picture(image_file=config['ReplacementImageDarkSmall'],
                                             left=pptx.util.Inches(logo_positions['bottom-right']['left']),
                                             top=pptx.util.Inches(logo_positions['bottom-right']['top'])
                                             )

            # print(check_for_logo_collision_updated(slide))
        else:  # For every other slide that is not a title slide
            # Check if old color scheme exists on the current slide's slide layout
            if _check_background_color(slide_layout=slide.slide_layout) is not None:
                # Add the Small, White BH logo to the bottom right of the slide
                slide.shapes.add_picture(image_file=config['ReplacementImageWhiteSmall'],
                                         left=pptx.util.Inches(logo_positions['bottom-right']['left']),
                                         top=pptx.util.Inches(logo_positions['bottom-right']['top'])
                                         )
            else:  # If old color scheme is not detected
                # Add the Small, Dark BH logo to the bottom right of the slide
                slide.shapes.add_picture(image_file=config['ReplacementImageDarkSmall'],
                                         left=pptx.util.Inches(logo_positions['bottom-right']['left']),
                                         top=pptx.util.Inches(logo_positions['bottom-right']['top'])
                                         )


def _delete_replacement_image(slide):
    """
    ---------------------------------- CURRENTLY NOT USED IN THE SCRIPT ----------------------------------
    Deletes the BH replacement logo from the current slide

    Parameters
    ----------
    slide: Slide object

    Returns None
    -------
    """

    for shape in slide.shapes:
        # For each shape on the slide check if it is a picture
        if 'pic' in shape.element.tag:
            # Search for the Non-Visual Drawing Properties of the picture
            element_pr = shape.element.find('.//ns0:cNvPr',
                                            {'ns0': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
            if element_pr.attrib is not None:
                if 'descr' in element_pr.attrib:
                    # if element_pr.attrib['descr'] == 'replacementImage2.png':
                    if element_pr.attrib['id'] == '1000' and element_pr.attrib['name'] == 'Picture 999':
                        # If the current element in the replacement image, delete it
                        pic = shape.element
                        pic_p = pic.getparent()
                        pic_p.remove(pic)


def check_for_logo_collision_updated(slide):
    """
    ---------------------------------- CURRENTLY NOT USED IN THE SCRIPT ----------------------------------
    Checks if the BH logo collides with other elements on the current slide

    Parameters
    ----------
    slide: Slide object

    Returns Boolean value indicating whether there is a collision
    -------
    """

    prefix_map = {
        'ns0': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'ns1': 'http://schemas.openxmlformats.org/drawingml/2006/main'
    }  # Contains the mapping from namespace mapping to full name for searching the XML tree

    img_element_pr = None  # Variable for holding image element properties
    other_shapes_pr = []  # List for storing properties of other shapes

    for shape in slide.shapes:
        if 'pic' in shape.element.tag:
            image_element = shape.element.find('.//ns0:cNvPr', prefix_map)
            if image_element.attrib is not None:
                if 'descr' in image_element.attrib:
                    #  if image_element.attrib['descr'] in ['replacementImage2.png', 'replacementImage.png']:
                    if image_element.attrib['id'] == '1000' and image_element.attrib['name'] == 'Picture 999':
                        img_element_pr = shape.element
                        # break
        shape_element = shape.element.find('.//ns0:spPr', prefix_map)  # Find the Shape properties element
        if shape_element is not None:
            other_shapes_pr.append(shape_element)

    img_element_pr_loc = get_element_location(img_element_pr, prefix_map)

    other_shapes_pr_loc = [get_element_location(shape, prefix_map) for shape in other_shapes_pr]

    for shape_pr_loc in other_shapes_pr_loc:
        if shape_pr_loc is not None:
            x_shape_start = shape_pr_loc['x']
            x_shape_end = (shape_pr_loc['x'] + shape_pr_loc['cx'])
            y_shape_start = shape_pr_loc['y']
            y_shape_end = (shape_pr_loc['y'] + shape_pr_loc['cy'])

            x_img_start = img_element_pr_loc['x']
            x_img_end = (img_element_pr_loc['x'] + img_element_pr_loc['cx'])
            y_img_start = img_element_pr_loc['y']
            y_img_end = (img_element_pr_loc['y'] + img_element_pr_loc['cy'])

            collision = False

            if x_shape_start < x_img_start < x_shape_end and y_shape_start < y_img_end:
                collision = True

            elif x_shape_start < x_img_start < x_img_end and y_shape_start < y_img_start < y_shape_end:
                collision = True

            return collision


def get_element_location(element_node, prefix_map):
    """
    ---------------------------------- CURRENTLY NOT USED IN THE SCRIPT ----------------------------------

    Get the location and size of the element on the slide

    Parameters
    ----------
    element_node: XML _Element object
    prefix_map: dictionary
        Mapping from namespace to full name

    Returns Dictionary containing the X, Y coordinates of the element on the slide and the width and height of the
    element --> {'x': int, 'y': int, 'cx': int, 'cy': int} -------
    """

    try:
        # Get the X,Y coordinates of the element on the slide stored in the xfrm/off _Element attributes
        # prefix_map = 'ns1': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        off = element_node.find('.//ns1:off', prefix_map).attrib
        # Get the CX, CY - how far the _Element extends on the x and y-axis, stored in the xfrm/ext attributes
        # prefix_map = 'ns1': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ext = element_node.find('.//ns1:ext', prefix_map).attrib

        # Convert the string values to int
        locations = {key: int(value) for (key, value) in {**off, **ext}.items()}

        return locations  # Return the locations dictionary
    except AttributeError:
        # If the current element is missing the 'off', 'ext' sub-elements return None
        return None


def disable_background_graphics(zip_in, zip_out, xml_path):
    """
    Disables the background graphics on the slide and adds copyright information

    Parameters
    ----------
    zip_in: input ZipFile object
    zip_out: output ZipFile object
    xml_path: string
        Path to the PowerPoint slide xml file

    Returns None
    -------
    """

    xml_tree = zip_in.open(xml_path).read()  # Open and read in the xml file from the path
    root = lxml.etree.fromstring(xml_tree)  # Parse the tree

    if root.attrib.get('showMasterSp') is None:  # If 'showMasterSp' attribute is not present in the xml file
        root.attrib['showMasterSp'] = '0'  # Set the showMasterSp attribute to 0 to hide background graphics

    xml_element = _create_copyright_element()

    sp_tree = root.find('.//ns0:spTree', {'ns0': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
    # sp_tree.insert(-1, xml_element)  # Append the xml element to the spTree
    sp_tree.append(xml_element)  # Append the xml element to the spTree

    # Write the updated xml file to the output file
    modified_xml = lxml.etree.tostring(root, encoding='utf-8', xml_declaration=True)
    zip_out.writestr(xml_path, modified_xml)


def convert_to_pdf_2(file_path, file_name, config):
    cv_pdf_2_docx = ConverterPdf2Docx(file_path)
    cv_pdf_2_docx.convert(config['InputFolder'] + f'/{file_name}.docx', start=0, end=None)
    cv_pdf_2_docx.close()


def process_file_pdf(file_in, file_out, config):
    file_path = file_in
    pdfs.append(file_path.split('\\')[-1])
    file_name = file_path.split('\\')[-1][0:-4]
    file_out_path = file_out.replace('.pdf', '.docx')

    if file_path.endswith('.pdf'):
        start_time = time()
        p = multiprocessing.Process(target=convert_to_pdf_2, args=(file_path, file_name, config))
        p.start()
        while time() - start_time < 60:
            if not p.is_alive():
                p.join()
                break
            sleep(.1)
        else:
            print('TIMED OUT - TERMINATING PROCESS')
            p.terminate()
            p.join()
            raise Exception('TIMED OUT - Process Terminated')

        """cv_pdf_2_docx = ConverterPdf2Docx(file_path)
        cv_pdf_2_docx.convert(config['InputFolder'] + f'/{file_name}.docx', start=0, end=None)
        cv_pdf_2_docx.close()"""

    """start_process_time = time()
    p2 = multiprocessing.Process(target=process_file_word, args=(
        config['InputFolder'] + f'/{file_name}.docx', file_out_path, config, header_images_pdf))
    p2.start()
    while time() - start_process_time < 120:
        if not p2.is_alive():
            p2.join()
            break
        sleep(.1)
    else:
        print('TIMED OUT - TERMINATING PROCESS')
        p2.terminate()
        p2.join()
        os.remove(config['InputFolder'] + f'/{file_name}.docx')
        raise Exception('TIMED OUT - Process Terminated')"""

    process_file_word(file_in=config['InputFolder'] + f'/{file_name}.docx',
                      file_out=file_out_path, config=config)

    os.remove(config['InputFolder'] + f'/{file_name}.docx')


def process_file(file_in, file_out, config):
    """
    Replaces the old BHGE logo and copyright information

    Parameters
    ----------
    file_in: string
        the path to the input file

    file_out: string
        the path to the output file

    config: dictionary
        content of the configuration file as dictionary with the tag as
        key and the information as value
    """
    duration = time()
    global file_run_times
    file_type = None
    file_extension = pathlib.Path(file_in).suffix
    match file_extension:
        case '.doc' | '.docx' | '.docm':
            process_file_word(file_in, file_out, config)
        case '.xlsx' | '.xls':
            process_file_excel(file_in, file_out, config)
            file_type = 'excel'
        case '.pptx':
            process_file_powerpoint(file_in, file_out, config)
            file_type = 'powerpoint'
        case '.pdf':
            process_file_pdf(file_in, file_out, config)
            file_type = 'pdf'

    if file_type is not None:
        file_run_times[file_type] = file_run_times[file_type] + (time() - duration)


def add_missing_images(zip_in, zip_out):
    """
    Adds missing images from one Zip archive to another.

    Parameters
    ----------
    zip_in : ZipFile
        The input ZipFile object from which to extract the image filenames.

    zip_out : ZipFile
        The output ZipFile object where the missing images will be added.
    """
    # Get a list of image filenames from zip_in
    image_files = [item.filename for item in zip_in.infolist() if '/media/' in item.filename]

    # Check if the image already exists in zip_out, and if not, add it
    for image_file in image_files:
        if image_file not in zip_out.namelist():
            image_data = zip_in.read(image_file)
            zip_out.writestr(image_file, image_data)


def convert_file(file, new_file_format):
    # Get the folder path and the base filename of the original file
    folder_path = os.path.dirname(file)
    base_name = os.path.basename(file)

    # Get file application depending on file type
    file_app = None
    extension = None
    opened_file = None
    if new_file_format == FILE_FORMAT_DOCX:
        file_app = win32.gencache.EnsureDispatch('Word.Application')
        extension = '.docx'

        # Open the .doc file
        opened_file = file_app.Documents.Open(file)

    elif new_file_format == FILE_FORMAT_XLSX:
        file_app = win32.gencache.EnsureDispatch('Excel.Application')
        extension = '.xlsx'

        # Open the .xls file
        opened_file = file_app.Workbooks.Open(file)

    # Construct the paths for the new file and the temporary new copy
    new_file = os.path.join(folder_path, os.path.splitext(base_name)[0] + extension)
    temp_new_file = os.path.join(folder_path, os.path.splitext(base_name)[0] + '_temp' + extension)

    # Save the document as the new format
    opened_file.SaveAs(temp_new_file, FileFormat=new_file_format)

    # Close the document and the application
    opened_file.Close()
    file_app.Quit()

    # Delete the original file
    os.remove(file)

    # Rename temp file to original name
    os.rename(temp_new_file, new_file)

    return new_file


def prepare_log(config):
    """
    Creates and prepares logfile

    Parameters
    ----------
    config: dictionary
        content of the configuration file as dictionary with the tag as
        key and the information as value

    Returns
    -------
    log: TextIOWrapper
        wrapper for the logfile
    """
    log_time = datetime.now().strftime("%Y-%m-%d_%H%M%S")
    log_name = f"BH_Rebrand_{log_time}.csv"
    log_path = os.path.join(config["LogFolder"], log_name)

    # Open Log - it will need to be closed in main function after processing all files
    log = open(log_path, "w+", encoding="utf-8")

    # Create the header for the log
    log.write("sep=;\n")
    log.write("Inputfile;Logo;NumLogos;Notes;LegacyText;Time;Warning\n")

    print(f"Logging in file {log_name}")

    return log, log_path


def main():
    """
    Starts processing of files and conversion to PDF
    """
    script_run_time = time()

    # Multiprocessing manager - currently not utilized in script - needed when file processing is done using separate
    # processes to share the header_images_pdf variable across processes
    # For regular file processing the header_images global variable is used by default
    # manager = multiprocessing.Manager()
    # header_images_pdf = manager.dict()

    if len(sys.argv) == 2:
        # Put elements of config file into a dictionary
        config = map_config(sys.argv[1])

        # Check if folders exist and create them if necessary
        for folder in FOLDERS:
            if not os.path.isdir(config[folder]):
                os.mkdir(config[folder])

        # Clean the project before starting (used for testing)
        delete_all_contents(config)

        # Prepare the log file
        log, log_path = prepare_log(config)

        process_file_failure_count = 0

        # PDF conversion is currently NOT required for file processing
        """pdf_conversion = False
        if "true" in config["PDF"].lower():
            pdf_conversion = True
            print("Starting COM Server for PDF conversion")
            # Start COM servers
            start_time = time()
            # Get a list of the filetypes that are going to be rebranded
            filetypes = get_filetypes_in_folder(config["InputFolder"])

            # Start the necessary COM Servers
            word, excel, power_point = start_com_servers(filetypes)

            # End timer and output time
            print(f"Starting server took {time() - start_time:.3f} seconds")"""

        # Check if input directory exists
        if os.path.isdir(config["InputFolder"]):
            # Decide on output folder
            if config["CompareLogoByPixels"]:
                output_folder = config["HeaderImageReplacedFoler"]
            else:
                output_folder = config["OutputFolder"]

            # Get total file count and set counter
            file_count = len(os.listdir(config['InputFolder']))
            file_counter = 1

            # Sort files to process Word files first before other file types
            sorted_files = sorted(os.listdir(config["InputFolder"]),
                                  key=lambda x: (not x.lower().endswith((".docx", ".doc", ".docm")), x))

            # Loop over every file in the directory
            for file in sorted_files:
                # Start timer
                start_time = time()

                # Create input and output path and start file processing
                file_in = os.path.join(config["InputFolder"], file)
                file_out = os.path.join(output_folder, file)

                # Get the document number for the current file being processed
                doc_num = file_in.split('\\')[-1]

                # If the document number is present in the 'ProcessedFolder', the file has been already re-branded
                if doc_num in os.listdir(config['ProcessedFolder']):
                    print(f'{file_in} skipped - File has already been processed')
                    log.write(f'{file_in};-;-;File has already been re-branded;File Skipped;-;-\n')
                    continue

                # Functionality for handling 'Saving as macro-enabled workbook' pop-up is not implemented in the
                # script as of right now skip .xlsm files
                if file_in.endswith(".xlsm"):
                    print(f'{file_in} skipped - XLSM format found')
                    log.write(f'{file_in};-;-;XLSM format found;File Skipped;-;-\n')
                    continue

                # Check if current file is a path to a folder
                if os.path.isdir(file_in):
                    # Ignore folder and continue
                    continue

                # Get the current filetype
                config = get_filetype(file, config)

                # Check if filetype is supported
                if "filetype" not in config:
                    print(f"Filetype of {file} not supported")
                    log.write(f"{file_in};-;-;Filetype not supported;-;-\n")
                    continue

                # Format the OldLogoPath
                config["OldLogoPath_formatted"] = config["OldLogoPath"].format(
                    filetype=config["filetype"])
                config["LegacyBHLogoPath_formatted"] = config["LegacyBHLogoPath"].format(
                    filetype=config["filetype"])

                # Process the file if it's not empty
                print(f"File {file_counter}/{file_count} -- Processing {file}")
                try:
                    if os.path.getsize(file_in) != 0:
                        log_information['Inputfile'] = file_in
                        process_file(file_in, file_out, config)
                        processing_time = time() - start_time
                        log_information['Time'] = str(processing_time)
                        print(f"Processing took {time() - start_time:.3f} seconds")
                        log.write(';'.join([str(value) for key, value in log_information.items() if key]) + '\n')
                    else:
                        print(f"File skipped because it's empty: {file}")
                        log.write(f"{file_in};-;-;File skipped because it's empty!;-;-\n")
                except Exception as e:
                    print(f'File failed to process! File name: {file}\nError: {e}')
                    process_file_failure_count += 1
                    log.write(f"{file_in};-;-;File failed to process!;Error:{e};-\n")
                # End timer and output time
                file_counter += 1
                reset_log_info()

                # Convert to PDF
                """if pdf_conversion:
                    # Start timer
                    start_converting_time = time()

                    # Convert file using COM server
                    print(f"Converting {file} to PDF")
                    config["PDFPath"] = os.path.join(output_folder,
                                                     file[0:file.rfind(".")] + ".pdf")

                    # Convert the file to pdf
                    convert_to_pdf(file_out, config, word, excel, power_point)

                    # End timer and output time
                    print(f"Converting took {time() - start_converting_time:.3f} seconds")"""
        else:
            print("Specify an existing input folder containing the documents and an output folder")

        # Loop through all the files, replacing any image inside the file's body that matches any logo in the catalog
        if config["CompareLogoByPixels"]:
            body_image_replace = time()
            file_counter = 1
            file_count = len(os.listdir(config['InputFolder']))
            body_replace_failure_count = 0
            global file_run_times

            # Sort files to process Word files first before other file types
            sorted_files = sorted(os.listdir(config["HeaderImageReplacedFoler"]),
                                  key=lambda x: (not x.lower().endswith((".docx", ".doc", ".docm")), x))

            # Loop over every file in the directory
            for file_body in sorted_files:
                process_time = time()

                # Create input and output path and start file processing
                file_in = os.path.join(config["HeaderImageReplacedFoler"], file_body)
                file_out = os.path.join(config["OutputFolder"], file_body)

                # Get the current filetype
                config = get_filetype(file_body, config)

                # Check if filetype is supported
                if "filetype" not in config:
                    print(f"Filetype of {file_body} not supported")
                    log.write(f"{file_in};-;-;Filetype not supported;-;-\n")
                    continue

                # Replace image inside body
                print(f'File {file_counter}/{file_count} -- Replacing body image for: {file_in}')
                try:
                    place_logo_body(file_in, file_out, config)
                    if file_in.endswith('.docx'):
                        log_information['Time'] = str(time() - process_time)
                        log.write(';'.join([str(value) for key, value in log_information.items() if key]) + '\n')
                        reset_log_info()
                except Exception as e:
                    print(f'Failed replacing the body images for file: {file_body} \nError: {str(e.__traceback__)}')
                    body_replace_failure_count += 1
                    log.write(f"{file_in};-;-;Failed replacing the body images!;Error:{e};-\n")

                file_counter += 1
                # Remove file from headerImageReplaced folder
                # os.remove(file_in)

                # If the current input file originally was 'pdf', convert it back to pdf
                if file_in.split('\\')[-1].replace('.docx', '.pdf') in pdfs:
                    file_out_pdf = file_out.replace('.docx', '.pdf')  # Update the output path with 'pdf' extension
                    try:
                        ConvertDocx2Pdf(file_out, file_out_pdf)  # Convert docx to pdf
                        os.remove(file_out)  # Remove the docx copy from the output folder
                    except Exception as e:
                        log.write(f'{file_out};-;-;Failed converting to PDF;Error{e};-\n')
                        continue

                    # Increment PDF processing timer
                    file_run_times['pdf'] = file_run_times['pdf'] + (time() - process_time)
                else:
                    # Increment word processing timer
                    file_run_times['word'] = file_run_times['word'] + (time() - process_time)

            print(
                f"All done! \nThe script ran for {strftime('%H:%M:%S', gmtime(time() - script_run_time))} seconds"
                f"\nReplacing the files' body images took: {strftime('%H:%M:%S', gmtime(time() - body_image_replace))} "
                f"seconds\nTotal image comparisons: {image_comparisons:,}\nProcessed file failed: "
                f"{process_file_failure_count}\nBody image replacement fails count:{body_replace_failure_count}\n\nTime "
                f"breakdown:\nWord: {format_time(file_run_times['word'])}\nExcel: "
                f"{format_time(file_run_times['excel'])}\nPowerPoint: {format_time(file_run_times['powerpoint'])}\nPDF: "
                f"{format_time(file_run_times['pdf'])}")

        # Close the log file
        log.close()

        # Close COM Server
        """if pdf_conversion:
            quit_com_servers(word, excel, power_point)"""

    else:
        print("Specify an existing config file")


def format_time(seconds):
    # Extract the integer and decimal parts of the seconds
    int_seconds = int(seconds)
    decimal_seconds = round((seconds - int_seconds) * 100)  # Round to two decimal places

    # Calculate hours, minutes, and seconds
    hours, remainder = divmod(int_seconds, 3600)
    minutes, seconds = divmod(remainder, 60)

    # Format the time as a string
    formatted_time = f"{hours:02d}:{minutes:02d}:{seconds:02d}.{decimal_seconds:02d} seconds"

    return formatted_time


def _modify_xml_image_crop_fit(zip_in, zip_out, xml_file_path):
    """
    Remove the rectangle element, from an XML file, that is responsible for cropping an image.
    """
    # Read the XML file from the ZipFile object
    with zip_in.open(xml_file_path) as xml_file:
        xml_data = xml_file.read()

    # Parse the XML data using LXML
    parser = lxml.etree.XMLParser(remove_blank_text=True)
    root = lxml.etree.fromstring(xml_data, parser=parser)

    # Namespace mapping for OpenXML drawing elements
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
    }

    # Find the <a:srcRect> element and remove it
    for src_rect_element in root.findall('.//a:srcRect', namespaces=ns):
        if src_rect_element is not None:
            parent = src_rect_element.getparent()
            parent.remove(src_rect_element)

    # Write the modified XML data to a temporary in-memory buffer
    modified_xml_data = lxml.etree.tostring(root, encoding='utf-8', xml_declaration=True)

    # Replace the original file inside the zip_in with the updated content
    zip_out.writestr(xml_file_path, modified_xml_data)

    return True


def compare_images(image_path1, image_path2):
    """
    Compare two images based on their pixel values and determine their similarity.

    Args:
        image_path1 (str): File path of the first image.
        image_path2 (str): File path of the second image.

    Returns:
        bool: True if the images are considered similar, False otherwise.
    """
    # Set output variable
    similarity = False
    global image_comparisons

    # Catch any unsupported image files
    try:
        # Open the images
        image1 = Image.open(image_path1)
        image2 = Image.open(image_path2)

        # Resize the images to ensure they have the same dimensions
        image1 = image1.resize(image2.size)

        # Convert the images to RGB mode (if they are not already)
        image1 = image1.convert("RGB")
        image2 = image2.convert("RGB")

        # Convert images to NumPy arrays
        image1_array = np.asarray(image1)
        image2_array = np.asarray(image2)

        # Calculate deviation
        deviation = np.mean(np.abs(image1_array - image2_array))

        # Pictures are similar if their deviation is lower than the set threshold
        similarity = deviation < 15

        # Increment counter
        image_comparisons += 1
    except Exception:
        pass
        # print('Image format not supported: ', str(e)) # FOR TESTING

    return similarity


def map_config(configfile):
    """
    This function creates a dictionary from the information gathered form
    the configuration file.

    Parameters
    ----------
    configfile: string
        the path to the configuration file

    Returns
    -------
    cfg_dict: dictionary
        content of the configuration file as dictionary with the tag as
        key and the information as value

    """
    cfg_dict = {"ReplaceString": [], "case-unsensitive": []}

    with open(configfile, encoding="utf-8") as conf:
        for line in conf:
            if not line.startswith("//"):
                if search(r'[\S\s]*? = [\S\s]*?', line) is not None:
                    (key, val) = line.split(' = ')
                    cfg_dict[key] = val.strip()
                elif search(r'[\S\s]*? -> [\S\s]*?', line) is not None:
                    line = str(line)
                    'line = bytes(line,  "utf-8")'
                    (old, new) = line.split(' -> ')
                    # pylint: disable=W1401
                    # Disable error from using \S\s in a binary string which is needed for regex
                    # if search(b'[\S\s]*? // case-unsensitive', new) is not None:
                    #    new = (new[0:new.find(b" // case-unsensitive")])
                    #    cfg_dict["case-unsensitive"].append(new)
                    cfg_dict["ReplaceString"].append([old.strip(), new.strip()])

    return cfg_dict


def delete_all_contents(config):
    """
    Delete all file contents from the specified folders.

    Args:
        config (dict): Dictionary containing the folder paths.

    Returns:
        None
    """
    folders = [
        config['BetweenFolder'],
        config['HeaderImageReplacedFoler'],
        # config['FoundLogosFolder'],
        # config['ImagesFolder'],
        config['OutputFolder']
    ]

    for folder in folders:
        _delete_folder_contents(folder)


def _delete_folder_contents(folder_path):
    """
    Delete all files within a folder.

    Args:
        folder_path (str): Path to the folder.

    Returns:
        None
    """
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        if os.path.isfile(file_path):
            os.remove(file_path)


def reset_log_info():
    """
    Reset log information to default values

    Returns:
        None
    """
    log_information['Inputfile'] = ''
    log_information['Logo'] = 'No Logo Found'
    log_information['NumLogos'] = 0
    log_information['Notes'] = ''
    log_information['LegacyText'] = ''
    log_information['Time'] = ''
    log_information['Warning'] = ''


if __name__ == "__main__":
    main()
